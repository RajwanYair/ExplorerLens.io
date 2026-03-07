# ExplorerLens.py — Document Decoder
# Copyright (c) 2026 ExplorerLens Project
"""
Renders thumbnails for PDF, DJVU, Office documents, and CHM files.
Uses PyMuPDF (fitz) for PDF/DJVU and python-docx/python-pptx for Office.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from PIL import Image

from .base import BaseDecoder

logger = logging.getLogger("explorerlens.decoders.document")

_DOC_EXTS = [
    ".pdf", ".djvu", ".chm",
    ".docx", ".pptx", ".xlsx",
    ".doc", ".ppt", ".xls",
    ".odt", ".odp",
]


class DocumentDecoder(BaseDecoder):
    """Handles PDF, DJVU, Office, and related document formats."""

    @property
    def name(self) -> str:
        return "DocumentDecoder"

    def supported_extensions(self) -> list[str]:
        return list(_DOC_EXTS)

    def decode(self, path: Path, size: int) -> Optional[Image.Image]:
        ext = path.suffix.lower()

        if ext in (".pdf", ".djvu"):
            return self._decode_mupdf(path, size)
        if ext == ".pptx":
            return self._decode_pptx(path, size)
        if ext == ".docx":
            return self._decode_docx(path, size)
        if ext == ".xlsx":
            return self._decode_xlsx(path, size)

        # Generic fallback — try PyMuPDF which handles many formats
        return self._decode_mupdf(path, size)

    # ── PyMuPDF (PDF / DJVU / XPS / EPUB) ───────────────────────────

    @staticmethod
    def _decode_mupdf(path: Path, size: int) -> Optional[Image.Image]:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            if doc.page_count == 0:
                doc.close()
                return None
            page = doc.load_page(0)
            # Calculate zoom to fit requested size
            rect = page.rect
            zoom = size / max(rect.width, rect.height)
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            doc.close()
            return img
        except ImportError:
            logger.warning("PyMuPDF not installed — skipping: %s", path)
            return None
        except Exception as exc:
            logger.debug("MuPDF decode failed for %s: %s", path, exc)
            return None

    # ── PPTX ─────────────────────────────────────────────────────────

    @staticmethod
    def _decode_pptx(path: Path, size: int) -> Optional[Image.Image]:
        """Render first slide of a PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Emu

            prs = Presentation(str(path))
            if not prs.slides:
                return None

            slide = prs.slides[0]
            # Try to find an image shape on the first slide
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    blob = shape.image.blob
                    return Image.open(io.BytesIO(blob))

            # No image found — generate a text-based placeholder
            return DocumentDecoder._text_placeholder(
                "PPTX", path.stem, size
            )
        except ImportError:
            logger.warning("python-pptx not installed — skipping: %s", path)
            return None
        except Exception as exc:
            logger.debug("PPTX decode failed for %s: %s", path, exc)
            return None

    # ── DOCX ─────────────────────────────────────────────────────────

    @staticmethod
    def _decode_docx(path: Path, size: int) -> Optional[Image.Image]:
        """Extract first image or generate text preview from DOCX."""
        try:
            from docx import Document

            doc = Document(str(path))

            # Try to find embedded images
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    blob = rel.target_part.blob
                    return Image.open(io.BytesIO(blob))

            # No image — generate text preview
            text_lines = []
            for para in doc.paragraphs[:10]:
                if para.text.strip():
                    text_lines.append(para.text.strip())
            preview = "\n".join(text_lines[:5]) or path.stem
            return DocumentDecoder._text_placeholder("DOCX", preview, size)

        except ImportError:
            logger.warning("python-docx not installed — skipping: %s", path)
            return None
        except Exception as exc:
            logger.debug("DOCX decode failed for %s: %s", path, exc)
            return None

    # ── XLSX ─────────────────────────────────────────────────────────

    @staticmethod
    def _decode_xlsx(path: Path, size: int) -> Optional[Image.Image]:
        """Generate a simple spreadsheet-icon thumbnail for XLSX."""
        return DocumentDecoder._text_placeholder("XLSX", path.stem, size)

    # ── Placeholder ──────────────────────────────────────────────────

    @staticmethod
    def _text_placeholder(fmt: str, text: str,
                          size: int) -> Image.Image:
        """Generate a document-style placeholder with format badge."""
        from PIL import ImageDraw, ImageFont

        colors = {
            "PDF": (200, 50, 50),
            "PPTX": (210, 100, 30),
            "DOCX": (40, 90, 200),
            "XLSX": (30, 140, 60),
        }
        accent = colors.get(fmt, (100, 100, 100))

        canvas = Image.new("RGB", (size, size), (250, 250, 250))
        draw = ImageDraw.Draw(canvas)

        # Header bar
        draw.rectangle([0, 0, size, size // 6], fill=accent)

        # Format badge
        try:
            font_big = ImageFont.truetype("segoeui.ttf", size // 8)
            font_sm = ImageFont.truetype("segoeui.ttf", size // 14)
        except Exception:
            font_big = ImageFont.load_default()
            font_sm = font_big

        draw.text((size // 20, size // 30), fmt,
                  fill=(255, 255, 255), font=font_big)

        # Truncate text preview
        preview = text[:60]
        draw.text((size // 20, size // 4), preview,
                  fill=(60, 60, 60), font=font_sm)

        return canvas
